import os
from PyPDF2 import PdfReader
from pptx import Presentation

DOCS_DIR = "AMIT_RAG_Chatbot\docs"

def read_pdf(path):
    text = ""
    reader = PdfReader(path)
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def read_pptx(path):
    text = ""
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def chunk_text(text, chunk_size=500, overlap=50):
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start += chunk_size - overlap
    return chunks

def ingest_docs():
    all_chunks = []

    docs_files = os.listdir(DOCS_DIR)

    for file in docs_files:
        path = os.path.join(DOCS_DIR, file)
        text = ""

        if file.endswith(".pdf"):
            text = read_pdf(path)
        elif file.endswith(".pptx"):
            text = read_pptx(path)

        if text:
            chunks = chunk_text(text)   
            all_chunks.extend(chunks)  

    print(f"✅ Loaded {len(all_chunks)} chunks from {len(docs_files)} files")
    return all_chunks

if __name__ == "__main__":
    chunks = ingest_docs()
    print(chunks[:2])
